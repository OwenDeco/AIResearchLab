from __future__ import annotations

import csv
import logging
import re

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def _table_to_markdown(table: list) -> str:
    """Convert a list-of-rows (each row a list of cell strings) to a Markdown table."""
    if not table:
        return ""
    rows = [[str(cell or "").strip() for cell in row] for row in table]
    rows = [r for r in rows if any(c for c in r)]
    if not rows:
        return ""
    col_count = max(len(r) for r in rows)
    rows = [r + [""] * (col_count - len(r)) for r in rows]
    header = "| " + " | ".join(rows[0]) + " |"
    separator = "| " + " | ".join(["---"] * col_count) + " |"
    body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
    return "\n".join(filter(None, [header, separator, body]))


def _clean_text(text: str) -> str:
    """Post-process extracted text: fix hyphenation, normalize whitespace."""
    if not text:
        return ""
    # Fix hyphenated line breaks ("informa-\ntion" → "information")
    text = re.sub(r"(\w)-\n(\w)", r"\1\2", text)
    # Collapse multiple spaces on a single line
    text = re.sub(r"[ \t]{2,}", " ", text)
    # Collapse 3+ blank lines to 2
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _is_scanned_page(text: str, min_chars: int = 50) -> bool:
    """Return True if the page likely contains no extractable text (scanned image)."""
    return len(text.strip()) < min_chars


def _ocr_image(image) -> str:
    """Run Tesseract OCR on a PIL image. Returns empty string if not available."""
    try:
        import pytesseract  # type: ignore
        return pytesseract.image_to_string(image)
    except ImportError:
        return ""
    except Exception as exc:
        logger.debug("OCR failed: %s", exc)
        return ""


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def _load_pdf(file_path: str) -> str:
    """
    Extract text from a PDF using pdfplumber with:
    - Tables extracted as Markdown (avoids garbled table text)
    - Text extracted from non-table regions only
    - OCR fallback for scanned pages (requires pytesseract + Pillow)
    - PyMuPDF fallback if pdfplumber is not installed
    """
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        logger.warning("pdfplumber not installed; falling back to PyMuPDF. Install with: pip install pdfplumber")
        return _load_pdf_pymupdf(file_path)

    pages_text: list[str] = []

    try:
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                page_parts: list[str] = []

                # Detect tables and record their bounding boxes
                try:
                    table_objects = page.find_tables()
                    table_bboxes = [t.bbox for t in table_objects]
                except Exception:
                    table_objects = []
                    table_bboxes = []

                # Extract text from areas outside tables
                if table_bboxes:
                    try:
                        non_table_page = page.filter(
                            lambda obj, bboxes=table_bboxes: not any(
                                obj.get("x0", 0) >= bbox[0] - 1
                                and obj.get("top", 0) >= bbox[1] - 1
                                and obj.get("x1", 0) <= bbox[2] + 1
                                and obj.get("bottom", 0) <= bbox[3] + 1
                                for bbox in bboxes
                            )
                        )
                        text = non_table_page.extract_text() or ""
                    except Exception:
                        text = page.extract_text() or ""
                else:
                    text = page.extract_text() or ""

                # OCR fallback for scanned pages
                if _is_scanned_page(text):
                    try:
                        img = page.to_image(resolution=200).original
                        ocr_text = _ocr_image(img)
                        if ocr_text.strip():
                            text = ocr_text
                            logger.info("OCR used for page %d in %s", page_num, file_path)
                    except Exception as exc:
                        logger.debug("OCR attempt failed on page %d: %s", page_num, exc)

                if text.strip():
                    page_parts.append(text.strip())

                # Append tables as Markdown
                for table_obj in table_objects:
                    try:
                        md = _table_to_markdown(table_obj.extract())
                        if md:
                            page_parts.append(md)
                    except Exception:
                        pass

                if page_parts:
                    pages_text.append("\n\n".join(page_parts))

    except Exception as exc:
        logger.error("pdfplumber extraction failed for %s: %s — trying PyMuPDF", file_path, exc)
        return _load_pdf_pymupdf(file_path)

    return _clean_text("\n\n".join(pages_text))


def _load_pdf_pymupdf(file_path: str) -> str:
    """Basic PyMuPDF text extraction — used as fallback."""
    try:
        import fitz  # type: ignore  # PyMuPDF

        doc = fitz.open(file_path)
        pages = [page.get_text("text") for page in doc]
        doc.close()
        return _clean_text("\n".join(p for p in pages if p))
    except Exception as exc:
        logger.error("PyMuPDF loading failed for %s: %s", file_path, exc)
        raise


# ---------------------------------------------------------------------------
# Plain text
# ---------------------------------------------------------------------------

def _load_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="replace") as fh:
        return fh.read()


# ---------------------------------------------------------------------------
# Markdown
# ---------------------------------------------------------------------------

_MD_CODE_BLOCK = re.compile(r"```.*?```", re.DOTALL)
_MD_INLINE_CODE = re.compile(r"`[^`]+`")
_MD_HEADING = re.compile(r"^#{1,6}\s+", re.MULTILINE)
_MD_BOLD_ITALIC = re.compile(r"(\*{1,3}|_{1,3})(.*?)\1")
_MD_LINK = re.compile(r"!?\[([^\]]*)\]\([^\)]*\)")
_MD_HTML_TAG = re.compile(r"<[^>]+>")
_MD_HORIZONTAL_RULE = re.compile(r"^[-*_]{3,}\s*$", re.MULTILINE)


def _strip_markdown(text: str) -> str:
    """Remove Markdown syntax, leaving mostly plain prose."""
    text = _MD_CODE_BLOCK.sub(" ", text)
    text = _MD_INLINE_CODE.sub(" ", text)
    text = _MD_HEADING.sub("", text)
    text = _MD_BOLD_ITALIC.sub(r"\2", text)
    text = _MD_LINK.sub(r"\1", text)
    text = _MD_HTML_TAG.sub(" ", text)
    text = _MD_HORIZONTAL_RULE.sub(" ", text)
    lines = [line.rstrip() for line in text.splitlines()]
    text = "\n".join(lines)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _load_md(file_path: str) -> str:
    return _strip_markdown(_load_txt(file_path))


# ---------------------------------------------------------------------------
# Word (.docx)
# ---------------------------------------------------------------------------

def _load_docx(file_path: str) -> str:
    """Extract text and tables from a Word document in document order."""
    try:
        from docx import Document as DocxDocument  # type: ignore
        from docx.oxml.ns import qn  # type: ignore
        from docx.table import Table  # type: ignore
        from docx.text.paragraph import Paragraph  # type: ignore
    except ImportError:
        raise RuntimeError("python-docx not installed. Run: pip install python-docx")

    doc = DocxDocument(file_path)
    parts: list[str] = []

    # Iterate body children in document order (preserves table/paragraph interleaving)
    for child in doc.element.body.iterchildren():
        if child.tag == qn("w:p"):
            text = Paragraph(child, doc).text.strip()
            if text:
                parts.append(text)
        elif child.tag == qn("w:tbl"):
            rows = [
                [cell.text.strip() for cell in row.cells]
                for row in Table(child, doc).rows
            ]
            md = _table_to_markdown(rows)
            if md:
                parts.append(md)

    return _clean_text("\n\n".join(parts))


# ---------------------------------------------------------------------------
# Excel (.xlsx)
# ---------------------------------------------------------------------------

def _load_xlsx(file_path: str) -> str:
    """Extract all sheets from an Excel workbook as Markdown tables."""
    try:
        import openpyxl  # type: ignore
    except ImportError:
        raise RuntimeError("openpyxl not installed. Run: pip install openpyxl")

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(cells)
        if rows:
            parts.append(f"## Sheet: {sheet_name}\n\n" + _table_to_markdown(rows))

    wb.close()
    return _clean_text("\n\n".join(parts))


# ---------------------------------------------------------------------------
# PowerPoint (.pptx)
# ---------------------------------------------------------------------------

def _load_pptx(file_path: str) -> str:
    """
    Extract text and tables from each slide of a PowerPoint presentation.

    Shapes are sorted by visual position (top→bottom, left→right) so that
    content reads in the same order a human would scan the slide.
    Speaker notes are appended after the slide body when present.
    """
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation(file_path)
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []

        # Sort shapes top→bottom, left→right by their position on the slide.
        # pptx stores positions in EMUs; we bucket by vertical band (10% of
        # slide height) so that shapes on roughly the same row sort left→right.
        slide_height = prs.slide_height or 1
        def _shape_sort_key(shape):
            try:
                row_band = round(shape.top / slide_height * 10)
                return (row_band, shape.left)
            except Exception:
                return (0, 0)

        sorted_shapes = sorted(slide.shapes, key=_shape_sort_key)

        for shape in sorted_shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
            if shape.has_table:
                rows = [
                    [cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows
                ]
                md = _table_to_markdown(rows)
                if md:
                    slide_parts.append(md)

        # Speaker notes — often contain the full explanation of the slide
        try:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text.strip()
            if notes_text:
                slide_parts.append(f"Notes: {notes_text}")
        except Exception:
            pass

        if slide_parts:
            parts.append(f"## Slide {i}\n\n" + "\n\n".join(slide_parts))

    return _clean_text("\n\n".join(parts))


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------

def _load_html(file_path: str) -> str:
    """Extract readable text from an HTML file."""
    raw = _load_txt(file_path)
    try:
        from bs4 import BeautifulSoup  # type: ignore
        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return _clean_text(soup.get_text(separator="\n"))
    except ImportError:
        # Fallback: strip tags with regex
        return _clean_text(re.sub(r"<[^>]+>", " ", raw))


# ---------------------------------------------------------------------------
# CSV
# ---------------------------------------------------------------------------

def _load_csv(file_path: str) -> str:
    """Extract CSV content as a Markdown table."""
    rows: list[list[str]] = []
    with open(file_path, "r", encoding="utf-8", errors="replace", newline="") as fh:
        reader = csv.reader(fh)
        for row in reader:
            if any(c.strip() for c in row):
                rows.append(row)
    return _clean_text(_table_to_markdown(rows))


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

_LOADERS = {
    "pdf":      _load_pdf,
    "txt":      _load_txt,
    "md":       _load_md,
    "markdown": _load_md,
    "docx":     _load_docx,
    "xlsx":     _load_xlsx,
    "pptx":     _load_pptx,
    "html":     _load_html,
    "htm":      _load_html,
    "csv":      _load_csv,
}

SUPPORTED_TYPES = set(_LOADERS.keys())


def load_document(file_path: str, file_type: str) -> str:
    """
    Load *file_path* and return its plain-text content.

    Supported types: pdf, txt, md, docx, xlsx, pptx, html, htm, csv.
    Falls back to plain-text read for unknown types.
    """
    file_type = file_type.lower().lstrip(".")
    loader = _LOADERS.get(file_type)
    if loader is None:
        logger.warning("Unknown file type '%s'; attempting plain-text read.", file_type)
        return _load_txt(file_path)
    return loader(file_path)
